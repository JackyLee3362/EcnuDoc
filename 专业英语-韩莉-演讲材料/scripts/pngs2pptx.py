from pptx import Presentation
from pptx.util import Inches

def conf_prs():
    prs = Presentation()
    prs.slide_height = Inches(9)
    prs.slide_width = Inches(16)
    return prs
def run(
    num=23, 
    ext='png',
    pdfimg_path='./pdfimg', 
    save_path='test.pptx'
    ):

    prs = conf_prs()
    blank_slide_layout = prs.slide_layouts[6]

    slide = [0] * num
    top = left = Inches(0)
    height = Inches(9)

    for i in range(num):
        img_path = fr'{pdfimg_path}/{i}.{ext}'
        slide[i] = prs.slides.add_slide(blank_slide_layout)
        slide[i].shapes.add_picture(img_path, left, top, height=height)
    
    prs.save(save_path)
if __name__ == '__main__':
    run()
    
    

    

    

    

    